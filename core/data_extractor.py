"""Extract chart data from PPTX files into pandas DataFrames."""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from io import BytesIO

import pandas as pd
from pptx import Presentation
from pptx.chart.chart import Chart
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn


XY_CHART_TYPES = {
    XL_CHART_TYPE.XY_SCATTER,
    XL_CHART_TYPE.XY_SCATTER_LINES,
    XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
    XL_CHART_TYPE.XY_SCATTER_SMOOTH,
    XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS,
}

CATEGORY_LABEL = "קטגוריה"


def _extract_xy_x_values(series, n: int) -> list:
    """X values for an XY (scatter) series.

    python-pptx only exposes Y values publicly (series.values reads c:yVal),
    so read c:xVal from the series XML directly. When the series has no
    explicit X values, PowerPoint plots against the point index — mirror that.
    """
    xVal = series._element.xVal
    if xVal is None:
        return list(range(1, n + 1))
    return [xVal.pt_v(idx) for idx in range(xVal.ptCount_val)]


def is_percentage_format(fmt: str) -> bool:
    if not fmt:
        return False
    cleaned = re.sub(r'"[^"]*"', '', fmt)
    cleaned = re.sub(r'\\\.', '', cleaned)
    return '%' in cleaned


def _extract_series_formats_by_index(chart: Chart) -> list[str]:
    formats = []
    chart_xml = chart.part._element
    for ser in chart_xml.iter(qn('c:ser')):
        fmt_code = "General"
        val = ser.find(qn('c:val'))
        if val is not None:
            num_ref = val.find(qn('c:numRef'))
            if num_ref is not None:
                num_cache = num_ref.find(qn('c:numCache'))
                if num_cache is not None:
                    fc = num_cache.find(qn('c:formatCode'))
                    if fc is not None and fc.text:
                        fmt_code = fc.text
        formats.append(fmt_code)
    return formats


@dataclass
class ChartInfo:
    slide_index: int
    shape_name: str
    chart_type: int
    dataframe: pd.DataFrame
    is_xy: bool = False
    series_names: list = field(default_factory=list)
    series_formats: dict = field(default_factory=dict)
    chart_title: str = ""
    shape_id: int = 0

    @property
    def key(self):
        return (self.slide_index, self.shape_name)


def _extract_chart_data(chart: Chart):
    chart_type = chart.chart_type
    is_xy = chart_type in XY_CHART_TYPES

    plot = chart.plots[0]
    series_list = list(plot.series)
    series_names = [s.name if s.name else f"סדרה {i+1}" for i, s in enumerate(series_list)]

    format_list = _extract_series_formats_by_index(chart)
    series_formats = {}
    for i, name in enumerate(series_names):
        if i < len(format_list):
            series_formats[name] = format_list[i]

    if is_xy:
        data = {}
        for i, series in enumerate(series_list):
            y_vals = list(series.values)
            x_vals = _extract_xy_x_values(series, len(y_vals))
            data[f"X_{series_names[i]}"] = x_vals
            data[f"Y_{series_names[i]}"] = y_vals
        # Series may have different point counts — pad so DataFrame accepts them
        max_len = max((len(v) for v in data.values()), default=0)
        for col, vals in data.items():
            data[col] = vals + [None] * (max_len - len(vals))
        display_df = pd.DataFrame(data)
    else:
        try:
            categories = [str(c) for c in plot.categories]
        except Exception:
            categories = [str(i + 1) for i in range(len(list(series_list[0].values)))]

        display_data = {CATEGORY_LABEL: categories}
        for i, series in enumerate(series_list):
            values = list(series.values)
            while len(values) < len(categories):
                values.append(None)
            values = values[:len(categories)]

            name = series_names[i]
            fmt = series_formats.get(name, "General")
            if is_percentage_format(fmt):
                display_data[name] = [
                    round(v * 100, 2) if v is not None else None
                    for v in values
                ]
            else:
                display_data[name] = values

        display_df = pd.DataFrame(display_data)

    return display_df, is_xy, series_names, series_formats


def extract_all_charts(pptx_bytes: bytes) -> list[ChartInfo]:
    prs = Presentation(BytesIO(pptx_bytes))
    charts = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_chart:
                continue

            chart = shape.chart
            try:
                display_df, is_xy, series_names, series_formats = _extract_chart_data(chart)
                title_text = ""
                if chart.has_title and chart.chart_title and chart.chart_title.has_text_frame:
                    title_text = chart.chart_title.text_frame.text.strip()
                info = ChartInfo(
                    slide_index=slide_idx,
                    shape_name=shape.name,
                    chart_type=chart.chart_type,
                    dataframe=display_df,
                    is_xy=is_xy,
                    series_names=series_names,
                    series_formats=series_formats,
                    chart_title=title_text,
                    shape_id=shape.shape_id,
                )
                charts.append(info)
            except Exception as e:
                print(f"Warning: Could not extract chart '{shape.name}' on slide {slide_idx + 1}: {e}")

    return charts
