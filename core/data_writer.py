"""Write edited DataFrames back into PPTX charts."""

from __future__ import annotations

from io import BytesIO

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.oxml.ns import qn
from lxml import etree
from openpyxl import load_workbook

from core.data_extractor import is_percentage_format


def _display_to_raw(df: pd.DataFrame, series_formats: dict) -> pd.DataFrame:
    raw_df = df.copy()
    for col in df.columns[1:]:
        fmt = series_formats.get(col, "General")
        if is_percentage_format(fmt):
            raw_df[col] = df[col].apply(
                lambda v: v / 100.0 if pd.notna(v) else None
            )
    return raw_df


def update_multiple_charts(
    pptx_bytes: bytes,
    updates: list,
) -> bytes:
    """Update multiple charts in a single parse/save cycle.

    Each update tuple: (slide_index, shape_name, display_df, is_xy, series_formats, shape_id)
    """
    prs = Presentation(BytesIO(pptx_bytes))

    for update in updates:
        slide_index, shape_name, display_df, is_xy, series_formats, _shape_id = update
        df = _display_to_raw(display_df, series_formats) if series_formats else display_df
        slide = prs.slides[slide_index]

        chart_shape = None
        for shape in slide.shapes:
            if shape.has_chart and shape.name == shape_name:
                chart_shape = shape
                break
        if chart_shape is None:
            continue

        chart = chart_shape.chart

        if is_xy:
            chart_data = XyChartData()
            col_names = df.columns.tolist()
            for i in range(0, len(col_names), 2):
                series_name = col_names[i].replace("X_", "")
                series = chart_data.add_series(series_name)
                x_vals = df.iloc[:, i].dropna().tolist()
                y_vals = df.iloc[:, i + 1].dropna().tolist()
                for x, y in zip(x_vals, y_vals):
                    series.add_data_point(x, y)
        else:
            chart_data = CategoryChartData()
            categories = df.iloc[:, 0].dropna().astype(str).tolist()
            chart_data.categories = categories
            for col in df.columns[1:]:
                values = df[col].tolist()
                values = [None if pd.isna(v) else float(v) for v in values]
                values = values[:len(categories)]
                while len(values) < len(categories):
                    values.append(None)
                chart_data.add_series(col, values)

        chart.replace_data(chart_data)
        if series_formats:
            _restore_format_codes(chart, series_formats)
            _format_embedded_excel(chart, series_formats)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def _format_embedded_excel(chart, series_formats: dict):
    try:
        xlsx_part = chart.part.chart_workbook.xlsx_part
        wb = load_workbook(BytesIO(xlsx_part.blob))
    except Exception:
        return

    ws = wb.active
    if ws is None or ws.max_row is None or ws.max_column is None:
        return

    format_values = list(series_formats.values())
    for col_offset, fmt_code in enumerate(format_values):
        excel_col = col_offset + 2
        if excel_col > ws.max_column:
            break
        for row in range(2, ws.max_row + 1):
            cell = ws.cell(row=row, column=excel_col)
            if cell.value is not None:
                cell.number_format = fmt_code

    buf = BytesIO()
    wb.save(buf)
    xlsx_part.blob = buf.getvalue()


def _restore_format_codes(chart, series_formats: dict):
    chart_xml = chart.part._element
    format_values = list(series_formats.values())

    for idx, ser in enumerate(chart_xml.iter(qn('c:ser'))):
        if idx >= len(format_values):
            break
        fmt_code = format_values[idx]
        val = ser.find(qn('c:val'))
        if val is not None:
            num_ref = val.find(qn('c:numRef'))
            if num_ref is not None:
                num_cache = num_ref.find(qn('c:numCache'))
                if num_cache is not None:
                    fc = num_cache.find(qn('c:formatCode'))
                    if fc is None:
                        fc = etree.SubElement(num_cache, qn('c:formatCode'))
                    fc.text = fmt_code
