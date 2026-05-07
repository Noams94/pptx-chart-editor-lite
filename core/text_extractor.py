"""Extract editable text from PPTX files (chart titles, placeholders, textboxes, table cells)."""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


# Hebrew labels for placeholder roles. Anything not listed falls back to
# "placeholder:<idx>".
_PLACEHOLDER_ROLE_LABELS = {
    PP_PLACEHOLDER.TITLE: "כותרת שקופית",
    PP_PLACEHOLDER.CENTER_TITLE: "כותרת שקופית",
    PP_PLACEHOLDER.SUBTITLE: "כותרת משנה",
    PP_PLACEHOLDER.BODY: "גוף תוכן",
    PP_PLACEHOLDER.OBJECT: "גוף תוכן",
    PP_PLACEHOLDER.HEADER: "כותרת עליונה",
    PP_PLACEHOLDER.FOOTER: "כותרת תחתונה",
    PP_PLACEHOLDER.DATE: "תאריך",
    PP_PLACEHOLDER.SLIDE_NUMBER: "מספר שקופית",
}

TITLE_ROLES = {"כותרת שקופית"}


@dataclass
class TextInfo:
    slide_index: int
    slide_title: str
    shape_id: int
    shape_name: str
    kind: str
    paragraph_index: int
    table_row: int  # -1 when not a table cell
    table_col: int  # -1 when not a table cell
    original_text: str
    # Position in EMU; used only for sorting. Defaults make group children fine.
    sort_top: int = 0
    sort_left: int = 0

    @property
    def location_label(self) -> str:
        if self.kind == "table_cell":
            if self.table_row == 0:
                return f"כותרת — עמודה {self.table_col + 1} פסקה {self.paragraph_index + 1}"
            return f"שורה {self.table_row + 1} עמודה {self.table_col + 1} פסקה {self.paragraph_index + 1}"
        return f"פסקה {self.paragraph_index + 1}"


def _paragraph_text(paragraph) -> str:
    return "".join(run.text for run in paragraph.runs)


def _safe_pos(shape, attr: str) -> int:
    try:
        v = getattr(shape, attr, None)
        return int(v) if v is not None else 0
    except Exception:
        return 0


def _placeholder_kind(shape) -> str:
    try:
        ph_type = shape.placeholder_format.type
        idx = shape.placeholder_format.idx
    except Exception:
        return "placeholder"
    label = _PLACEHOLDER_ROLE_LABELS.get(ph_type)
    if label:
        return label
    return f"placeholder:{idx}"


def _emit_text_frame(text_frame, slide_index, shape_id, shape_name, kind,
                     sort_top, sort_left, table_row=-1, table_col=-1):
    items = []
    for p_idx, paragraph in enumerate(text_frame.paragraphs):
        text = _paragraph_text(paragraph)
        if not text.strip():
            continue
        items.append(TextInfo(
            slide_index=slide_index,
            slide_title="",  # filled in later, once per slide
            shape_id=shape_id,
            shape_name=shape_name,
            kind=kind,
            paragraph_index=p_idx,
            table_row=table_row,
            table_col=table_col,
            original_text=text,
            sort_top=sort_top,
            sort_left=sort_left,
        ))
    return items


def _walk_shape(shape, slide_index: int, slide_height: int, results: list,
                inherited_top: int = None, inherited_left: int = None):
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                # Pass through positions of child shapes themselves; if a child
                # has no position, fall back to the group's.
                _walk_shape(
                    sub, slide_index, slide_height, results,
                    inherited_top=_safe_pos(shape, "top"),
                    inherited_left=_safe_pos(shape, "left"),
                )
            return
    except Exception:
        pass

    shape_id = getattr(shape, "shape_id", 0)
    shape_name = getattr(shape, "name", "")
    sort_top = _safe_pos(shape, "top") or (inherited_top or 0)
    sort_left = _safe_pos(shape, "left") or (inherited_left or 0)

    if getattr(shape, "has_chart", False):
        chart = shape.chart
        if chart.has_title and chart.chart_title and chart.chart_title.has_text_frame:
            results.extend(_emit_text_frame(
                chart.chart_title.text_frame,
                slide_index, shape_id, shape_name, "chart_title",
                sort_top, sort_left,
            ))
        return

    if getattr(shape, "has_table", False):
        table = shape.table
        for r_idx, row in enumerate(table.rows):
            for c_idx, cell in enumerate(row.cells):
                results.extend(_emit_text_frame(
                    cell.text_frame,
                    slide_index, shape_id, shape_name, "table_cell",
                    sort_top, sort_left,
                    table_row=r_idx, table_col=c_idx,
                ))
        return

    if getattr(shape, "has_text_frame", False):
        if getattr(shape, "is_placeholder", False):
            kind = _placeholder_kind(shape)
        else:
            # textbox: classify as title if in top quarter of slide
            if slide_height and sort_top and sort_top < slide_height * 0.25:
                kind = "כותרת תיבה"
            else:
                kind = "תיבת טקסט"
        results.extend(_emit_text_frame(
            shape.text_frame,
            slide_index, shape_id, shape_name, kind,
            sort_top, sort_left,
        ))


def _sort_key(t: TextInfo):
    # Top ascending, left descending (RTL), then stable tie-breakers
    return (
        t.sort_top,
        -t.sort_left,
        t.shape_id,
        t.table_row,
        t.table_col,
        t.paragraph_index,
    )


def _compute_slide_titles(by_slide: dict[int, list[TextInfo]]) -> dict[int, str]:
    titles = {}
    for slide_idx, items in by_slide.items():
        title = ""
        # Prefer an explicit TITLE/CENTER_TITLE placeholder
        for t in items:
            if t.kind in TITLE_ROLES:
                title = t.original_text
                break
        # Fallback: first item in visual order
        if not title and items:
            title = items[0].original_text
        titles[slide_idx] = title
    return titles


def extract_all_texts(pptx_bytes: bytes) -> list[TextInfo]:
    prs = Presentation(BytesIO(pptx_bytes))
    slide_height = int(prs.slide_height) if prs.slide_height else 0

    raw: list[TextInfo] = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            try:
                _walk_shape(shape, slide_idx, slide_height, raw)
            except Exception as e:
                print(f"Warning: Could not extract text from shape on slide {slide_idx + 1}: {e}")

    # Group by slide, sort each group visually, then assemble
    by_slide: dict[int, list[TextInfo]] = {}
    for t in raw:
        by_slide.setdefault(t.slide_index, []).append(t)
    for items in by_slide.values():
        items.sort(key=_sort_key)

    titles = _compute_slide_titles(by_slide)
    results: list[TextInfo] = []
    for slide_idx in sorted(by_slide):
        title = titles.get(slide_idx, "")
        for t in by_slide[slide_idx]:
            t.slide_title = title
            results.append(t)
    return results
