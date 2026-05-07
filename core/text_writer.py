"""Apply edited text values back into PPTX shapes, preserving first-run formatting."""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


def _iter_all_shapes(shapes):
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_all_shapes(shape.shapes)
                continue
        except Exception:
            pass
        yield shape


def _find_shape(slide, shape_id: int, shape_name: str):
    by_name = None
    for shape in _iter_all_shapes(slide.shapes):
        if getattr(shape, "shape_id", None) == shape_id:
            return shape
        if by_name is None and getattr(shape, "name", None) == shape_name:
            by_name = shape
    return by_name


def _resolve_text_frame(shape, kind: str, table_row: int, table_col: int):
    if kind == "chart_title":
        if not getattr(shape, "has_chart", False):
            return None
        chart = shape.chart
        if not chart.has_title:
            return None
        title = chart.chart_title
        if not title.has_text_frame:
            return None
        return title.text_frame
    if kind == "table_cell":
        if not getattr(shape, "has_table", False):
            return None
        try:
            return shape.table.cell(table_row, table_col).text_frame
        except Exception:
            return None
    # All other kinds (placeholder roles like "כותרת שקופית", textbox titles,
    # plain textboxes, etc.) live on the shape's own text_frame.
    if getattr(shape, "has_text_frame", False):
        return shape.text_frame
    return None


def _set_paragraph_text_preserving_format(paragraph, new_text: str):
    """Replace the paragraph's text while keeping the formatting of the first run."""
    p_elem = paragraph._p
    runs = list(p_elem.findall(qn('a:r')))

    if runs:
        first_run = runs[0]
        # Remove additional runs in this paragraph
        for extra in runs[1:]:
            p_elem.remove(extra)
        t = first_run.find(qn('a:t'))
        if t is None:
            from lxml import etree
            t = etree.SubElement(first_run, qn('a:t'))
        t.text = new_text
    else:
        # No runs — clear paragraph and add a fresh run via the high-level API
        # Drop existing inline children except paragraph properties
        for child in list(p_elem):
            if child.tag != qn('a:pPr'):
                p_elem.remove(child)
        run = paragraph.add_run()
        run.text = new_text


def update_multiple_texts(pptx_bytes: bytes, updates: list) -> bytes:
    """Apply text updates.

    Each update tuple:
        (slide_index, shape_id, shape_name, kind, paragraph_index,
         table_row, table_col, new_text)
    """
    prs = Presentation(BytesIO(pptx_bytes))

    for upd in updates:
        (slide_index, shape_id, shape_name, kind, paragraph_index,
         table_row, table_col, new_text) = upd

        try:
            slide = prs.slides[slide_index]
        except IndexError:
            continue

        shape = _find_shape(slide, shape_id, shape_name)
        if shape is None:
            continue

        text_frame = _resolve_text_frame(shape, kind, table_row, table_col)
        if text_frame is None:
            continue

        paragraphs = list(text_frame.paragraphs)
        if paragraph_index < 0 or paragraph_index >= len(paragraphs):
            continue

        _set_paragraph_text_preserving_format(paragraphs[paragraph_index], new_text)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()
